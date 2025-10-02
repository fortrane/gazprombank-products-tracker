import asyncio
from typing import Dict, Any
from datetime import date
import uuid
from pptx import Presentation
from matplotlib import pyplot as plt
import io

# placeholder LLM call
async def generate_report_text(data: Dict[str, Any]) -> str:
    """Сюда интегрируем LLM (HF) для создания текста отчета"""
    # заглушка
    return f"Executive Summary for {data['product']} from {data['from_date']} to {data['to_date']}"

async def enqueue_report_job(report_request: Dict[str, Any]) -> str:
    """Очередь на генерацию отчета, возвращает job_id"""
    job_id = str(uuid.uuid4())
    # Можно тут добавить Celery/BackgroundTasks для реальной обработки
    asyncio.create_task(generate_report(report_request, job_id))
    return job_id

async def generate_report(report_request: Dict[str, Any], job_id: str):
    """Сбор данных, LLM генерация текста, построение PPTX"""
    product = report_request['product']
    from_date = report_request['from_date']
    to_date = report_request['to_date']
    # placeholder: получить агрегаты из DB
    # db_client = report_request.get('db_client')
    # sentiment_data = await db_client.get_sentiment_aggregation(...)
    sentiment_data = [0.4, 0.2, 0.4]  # заглушка

    # 1) Генерация текста отчета через LLM
    report_text = await generate_report_text(report_request)

    # 2) Создание PPTX
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Report: {product}"
    slide.placeholders[1].text = report_text

    # Chart slide
    slide_chart = prs.slides.add_slide(prs.slide_layouts[5])
    fig, ax = plt.subplots()
    ax.bar(['Positive','Neutral','Negative'], sentiment_data)
    ax.set_title('Sentiment Distribution')
    buf = io.BytesIO()
    plt.savefig(buf, format='png')
    buf.seek(0)
    slide_chart.shapes.add_picture(buf, 100, 100, width=400, height=300)

    # Сохраняем файл локально или в S3/MinIO
    filename = f"report_{job_id}.pptx"
    prs.save(filename)
    print(f"Report saved: {filename}")
    # Можно сюда вернуть ссылку на S3
